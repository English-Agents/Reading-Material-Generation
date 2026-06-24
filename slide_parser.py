"""Parse Google Slides URLs or local PPTX/PDF files into structured ParsedSlide objects."""
from __future__ import annotations

import base64
import hashlib
import io
import re
from pathlib import Path
from typing import Optional

from pydantic import BaseModel


class EmbeddedImage(BaseModel):
    md5: str
    base64_data: str
    mime_type: str
    slide_index: int
    image_index: int


class ParsedSlide(BaseModel):
    slide_index: int
    title: Optional[str]
    body_text: str
    speaker_notes: str
    embedded_images: list[EmbeddedImage]
    source_url: str


_GOOGLE_SLIDES_RE = re.compile(
    r"https://docs\.google\.com/presentation/d/([A-Za-z0-9_-]+)"
)


def parse(source: str) -> list[ParsedSlide]:
    """Accept a Google Slides URL or a local PPTX/PDF path; return parsed slides."""
    if _GOOGLE_SLIDES_RE.match(source):
        return _parse_google_slides(source)
    path = Path(source)
    if path.suffix.lower() == ".pdf":
        return _parse_pdf(path)
    return _parse_pptx(path)


def _parse_pptx(path: Path) -> list[ParsedSlide]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER  # type: ignore[attr-defined]

    prs = Presentation(str(path))
    slides: list[ParsedSlide] = []

    # Placeholder types that represent a slide title in PowerPoint.
    _TITLE_PH_TYPES = {
        PP_PLACEHOLDER.TITLE,
        PP_PLACEHOLDER.CENTER_TITLE,
        PP_PLACEHOLDER.SUBTITLE,
    }

    def _is_title_placeholder(shape) -> bool:
        try:
            pf = shape.placeholder_format
            return pf is not None and (pf.idx == 0 or pf.type in _TITLE_PH_TYPES)
        except (ValueError, AttributeError):
            return False

    def _walk(shapes, idx, *, title_ref, text_parts, text_shapes, images, img_counter):
        """Recursively walk shapes — descends into groups and reads tables — so
        text inside grouped/templated shapes is not lost (the common reason a
        polished deck parses with zero titles)."""
        for shape in shapes:
            # Group → recurse into its children
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _walk(
                    shape.shapes, idx,
                    title_ref=title_ref, text_parts=text_parts,
                    text_shapes=text_shapes, images=images, img_counter=img_counter,
                )
                continue

            # Table → join all cell text
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        rows.append(" | ".join(cells))
                tbl_text = "\n".join(rows).strip()
                if tbl_text:
                    text_parts.append(tbl_text)
                    text_shapes.append(tbl_text)
                continue

            # Text frame
            if shape.has_text_frame:
                shape_text = shape.text_frame.text.strip()
                if shape_text:
                    if _is_title_placeholder(shape) and title_ref[0] is None:
                        title_ref[0] = shape_text
                    else:
                        text_parts.append(shape_text)
                    text_shapes.append(shape_text)

            # Picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    content_type = shape.image.content_type or "image/png"
                    images.append(
                        EmbeddedImage(
                            md5=hashlib.md5(blob).hexdigest(),
                            base64_data=base64.b64encode(blob).decode(),
                            mime_type=content_type,
                            slide_index=idx,
                            image_index=img_counter[0],
                        )
                    )
                    img_counter[0] += 1
                except Exception:
                    pass

    for idx, slide in enumerate(prs.slides):
        title_ref: list[Optional[str]] = [None]
        text_parts: list[str] = []
        text_shapes: list[str] = []   # all non-empty text, in document order — fallback source
        images: list[EmbeddedImage] = []
        img_counter = [0]

        _walk(
            slide.shapes, idx,
            title_ref=title_ref, text_parts=text_parts,
            text_shapes=text_shapes, images=images, img_counter=img_counter,
        )
        title = title_ref[0]

        # Fallback: no formal title placeholder (common in decks made from text
        # boxes, Google Slides, Canva, grouped templates). Use the first non-empty
        # text shape's first line as the title so the LLM receives a real topic.
        if title is None and text_shapes:
            first_line = text_shapes[0].splitlines()[0].strip()
            if first_line and len(first_line) <= 120:
                title = first_line
                if text_parts and text_parts[0].splitlines()[0].strip() == first_line:
                    remainder = "\n".join(text_parts[0].splitlines()[1:]).strip()
                    if remainder:
                        text_parts[0] = remainder
                    else:
                        text_parts.pop(0)

        # Speaker notes
        notes_text = ""
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip() if notes_tf else ""

        slides.append(
            ParsedSlide(
                slide_index=idx,
                title=title,
                body_text=" ".join(text_parts),
                speaker_notes=notes_text,
                embedded_images=images,
                source_url=str(path),
            )
        )

    return slides


def _parse_pdf(path: Path) -> list[ParsedSlide]:
    import pypandoc  # type: ignore[import-untyped]

    plain = pypandoc.convert_file(str(path), "plain")
    pages = plain.split("\f")  # form-feed character as page separator fallback
    if len(pages) == 1:
        # Split on double newlines as a heuristic if no form-feeds
        pages = [p for p in plain.split("\n\n\n") if p.strip()]

    slides: list[ParsedSlide] = []
    for idx, page in enumerate(pages):
        lines = [l.strip() for l in page.strip().splitlines() if l.strip()]
        title = lines[0] if lines else None
        body = " ".join(lines[1:]) if len(lines) > 1 else ""
        slides.append(
            ParsedSlide(
                slide_index=idx,
                title=title,
                body_text=body,
                speaker_notes="",
                embedded_images=[],
                source_url=str(path),
            )
        )
    return slides


def _parse_google_slides(url: str) -> list[ParsedSlide]:
    import base64 as _b64
    import json

    from google.oauth2 import service_account  # type: ignore[import-untyped]
    from googleapiclient.discovery import build  # type: ignore[import-untyped]

    from ppt_agent.config.settings import settings

    sa_json = json.loads(_b64.b64decode(settings.google_service_account_json).decode())
    creds = service_account.Credentials.from_service_account_info(
        sa_json,
        scopes=["https://www.googleapis.com/auth/drive.readonly"],
    )

    match = _GOOGLE_SLIDES_RE.match(url)
    file_id = match.group(1)  # type: ignore[union-attr]

    service = build("drive", "v3", credentials=creds)
    resp = (
        service.files()
        .export(
            fileId=file_id,
            mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        .execute()
    )

    tmp = io.BytesIO(resp)
    tmp.name = "presentation.pptx"

    # Write to a temp path so python-pptx can open it
    import tempfile

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(resp)
        tmp_path = Path(f.name)

    try:
        slides = _parse_pptx(tmp_path)
        for s in slides:
            object.__setattr__(s, "source_url", url)
        return slides
    finally:
        tmp_path.unlink(missing_ok=True)
