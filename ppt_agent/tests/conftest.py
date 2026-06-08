"""Shared pytest fixtures for ppt_agent tests."""
from __future__ import annotations

import base64
import io

import pytest


@pytest.fixture
def sample_pptx_bytes() -> bytes:
    """Build a minimal two-slide PPTX in memory for parser tests."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    # Slide 0 — title slide
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = "Introduction"
    s0.placeholders[1].text = "Subtitle text"

    # Slide 1 — concept slide with body text
    s1 = prs.slides.add_slide(blank_layout)
    from pptx.util import Pt
    txBox = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
    txBox.text_frame.text = "This slide explains key concepts about machine learning."

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_path(tmp_path, sample_pptx_bytes):
    """Write the sample PPTX to a temp file and return its path."""
    p = tmp_path / "test.pptx"
    p.write_bytes(sample_pptx_bytes)
    return p
